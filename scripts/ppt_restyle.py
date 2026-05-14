import sys
import os
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

NSMAP = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
}

SHAPE_TAGS = {'sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'}


def _find_title_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None


def _find_content_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def _find_title_shape(slide):
    ph = _find_title_placeholder(slide)
    if ph:
        return ph
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.lower().startswith('title'):
            return shape
    text_boxes = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        and s.has_text_frame
        and s.text_frame.text.strip()
    ]
    if text_boxes:
        return min(text_boxes, key=lambda s: (s.top, s.left))
    return None


def validate_template(prs):
    if len(prs.slides) != 2:
        raise ValueError(f"模板必须恰好有 2 页幻灯片，当前有 {len(prs.slides)} 页")

    cover = prs.slides[0]
    if not _find_title_placeholder(cover):
        raise ValueError("模板封面页（第1页）缺少标题占位符")

    content = prs.slides[1]
    if not _find_title_placeholder(content):
        raise ValueError("模板内容页（第2页）缺少标题占位符")
    if not _find_content_placeholder(content):
        raise ValueError("模板内容页（第2页）缺少内容占位符")


def _tag_local(element):
    tag = element.tag
    return tag.split('}')[-1] if '}' in tag else tag


def _remap_rids(element, rId_map):
    r_ns = NSMAP['r']
    attrs = [f'{{{r_ns}}}{a}' for a in ('embed', 'link', 'id', 'dm', 'lo', 'qs', 'cs')]
    for el in element.iter():
        for attr in attrs:
            val = el.get(attr)
            if val and val in rId_map:
                el.set(attr, rId_map[val])


def _copy_background(src_slide, dst_slide):
    src_cSld = src_slide._element.find('p:cSld', NSMAP)
    dst_cSld = dst_slide._element.find('p:cSld', NSMAP)
    if src_cSld is None or dst_cSld is None:
        return
    src_bg = src_cSld.find('p:bg', NSMAP)
    if src_bg is not None:
        dst_bg = dst_cSld.find('p:bg', NSMAP)
        if dst_bg is not None:
            dst_cSld.remove(dst_bg)
        dst_cSld.insert(0, deepcopy(src_bg))


def duplicate_slide(prs, slide_index):
    src_slide = prs.slides[slide_index]
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    dst_spTree = new_slide.shapes._spTree

    for child in list(dst_spTree):
        if _tag_local(child) in SHAPE_TAGS:
            dst_spTree.remove(child)

    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if not rel.is_external:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    for child in src_slide.shapes._spTree:
        if _tag_local(child) in SHAPE_TAGS:
            new_child = deepcopy(child)
            _remap_rids(new_child, rId_map)
            dst_spTree.append(new_child)

    _copy_background(src_slide, new_slide)
    return new_slide


def remove_slide(prs, slide_index):
    rId_lst = prs.slides._sldIdLst
    sldId = rId_lst[slide_index]
    rId = sldId.get(f'{{{NSMAP["r"]}}}id')
    prs.part.drop_rel(rId)
    rId_lst.remove(sldId)


def get_title_text(slide):
    shape = _find_title_shape(slide)
    if shape and shape.has_text_frame:
        return shape.text_frame.text
    return ""


def set_title_text(slide, text):
    ph = _find_title_placeholder(slide)
    if not ph or not ph.has_text_frame:
        return
    tf = ph.text_frame
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
            for run in para.runs[1:]:
                run._r.getparent().remove(run._r)
        else:
            para.text = text
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
    else:
        tf.text = text


def get_content_shapes(slide):
    title_shape = _find_title_shape(slide)
    title_el = title_shape._element if title_shape else None
    header_bottom = (title_shape.top + title_shape.height + Emu(150000)) if title_shape else 0

    content = []
    for shape in slide.shapes:
        if shape._element is title_el:
            continue
        if shape.top < 0:
            continue
        if header_bottom and (shape.top + shape.height) <= header_bottom:
            continue
        try:
            ph_fmt = shape.placeholder_format
        except ValueError:
            ph_fmt = None
        if ph_fmt is not None:
            if ph_fmt.idx in (0, 10, 11, 12):
                continue
        content.append(shape)
    return content


def compute_bounding_box(shapes):
    if not shapes:
        return None
    left = min(s.left for s in shapes)
    top = min(s.top for s in shapes)
    right = max(s.left + s.width for s in shapes)
    bottom = max(s.top + s.height for s in shapes)
    return (left, top, right, bottom)


def compute_mapping(src_bounds, dst_bounds):
    src_left, src_top, src_right, src_bottom = src_bounds
    dst_left, dst_top, dst_right, dst_bottom = dst_bounds

    src_w = src_right - src_left
    src_h = src_bottom - src_top
    dst_w = dst_right - dst_left
    dst_h = dst_bottom - dst_top

    if src_w == 0 or src_h == 0:
        return (1.0, dst_left, dst_top, src_left, src_top)

    scale = min(dst_w / src_w, dst_h / src_h)
    return (scale, dst_left, dst_top, src_left, src_top)


def map_position(left, top, width, height, mapping):
    scale, dst_left, dst_top, src_left, src_top = mapping
    new_left = int((left - src_left) * scale + dst_left)
    new_top = int((top - src_top) * scale + dst_top)
    new_width = int(width * scale)
    new_height = int(height * scale)
    return (new_left, new_top, new_width, new_height)


def get_content_area_bounds(slide):
    ph = _find_content_placeholder(slide)
    if ph:
        return (ph.left, ph.top, ph.left + ph.width, ph.top + ph.height)
    raise ValueError("无法确定内容区域边界")


def _get_mc_elements(slide, header_bottom):
    spTree = slide._element.find('p:cSld/p:spTree', NSMAP)
    if spTree is None:
        return []
    p_ns = NSMAP['p']
    results = []
    for mc in spTree.findall('mc:AlternateContent', NSMAP):
        choice = mc.find('mc:Choice', NSMAP)
        if choice is None:
            continue
        shape_el = None
        for tag_local in SHAPE_TAGS:
            shape_el = choice.find(f'{{{p_ns}}}{tag_local}')
            if shape_el is not None:
                break
        if shape_el is None:
            continue
        xfrm = shape_el.find('.//a:xfrm', NSMAP)
        if xfrm is None:
            xfrm = shape_el.find('p:xfrm', NSMAP)
        if xfrm is None:
            continue
        off = xfrm.find('a:off', NSMAP)
        ext = xfrm.find('a:ext', NSMAP)
        if off is None or ext is None:
            continue
        left = int(off.get('x', 0))
        top = int(off.get('y', 0))
        width = int(ext.get('cx', 0))
        height = int(ext.get('cy', 0))
        if top < 0:
            continue
        if header_bottom and (top + height) <= header_bottom:
            continue
        results.append((mc, left, top, width, height))
    return results


def _migrate_mc_element(dst_slide, mc_el, left, top, width, height, mapping, src_slide):
    new_mc = deepcopy(mc_el)
    new_left, new_top, new_width, new_height = map_position(
        left, top, width, height, mapping
    )
    for xfrm in new_mc.iter(f'{{{NSMAP["a"]}}}xfrm'):
        off = xfrm.find('a:off', NSMAP)
        if off is not None:
            off.set('x', str(new_left))
            off.set('y', str(new_top))
        ext = xfrm.find('a:ext', NSMAP)
        if ext is not None:
            ext.set('cx', str(new_width))
            ext.set('cy', str(new_height))

    r_ns = NSMAP['r']
    r_attrs = [f'{{{r_ns}}}{a}' for a in ('embed', 'link', 'id', 'dm', 'lo', 'qs', 'cs')]
    rId_map = {}
    for el in new_mc.iter():
        for attr in r_attrs:
            val = el.get(attr)
            if val and val not in rId_map and val in src_slide.part.rels:
                rel = src_slide.part.rels[val]
                if rel.is_external:
                    new_rId = dst_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[val] = new_rId
    if rId_map:
        _remap_rids(new_mc, rId_map)

    dst_slide.shapes._spTree.append(new_mc)


def migrate_content(dst_slide, src_slide, dst_bounds):
    content_shapes = get_content_shapes(src_slide)
    title_shape = _find_title_shape(src_slide)
    header_bottom = (title_shape.top + title_shape.height + Emu(150000)) if title_shape else 0
    mc_elements = _get_mc_elements(src_slide, header_bottom)

    if not content_shapes and not mc_elements:
        return

    content_ph = _find_content_placeholder(dst_slide)
    if content_ph:
        content_ph._element.getparent().remove(content_ph._element)

    src_bounds = compute_bounding_box(content_shapes)
    for _, ml, mt, mw, mh in mc_elements:
        mc_right, mc_bottom = ml + mw, mt + mh
        if src_bounds is None:
            src_bounds = (ml, mt, mc_right, mc_bottom)
        else:
            src_bounds = (
                min(src_bounds[0], ml), min(src_bounds[1], mt),
                max(src_bounds[2], mc_right), max(src_bounds[3], mc_bottom),
            )
    if src_bounds is None:
        return

    mapping = compute_mapping(src_bounds, dst_bounds)

    for shape in content_shapes:
        _migrate_single_shape(dst_slide, shape, mapping, src_slide)

    for mc_el, ml, mt, mw, mh in mc_elements:
        _migrate_mc_element(dst_slide, mc_el, ml, mt, mw, mh, mapping, src_slide)


def _migrate_single_shape(dst_slide, shape, mapping, src_slide):
    if hasattr(shape, 'has_chart') and shape.has_chart:
        _migrate_chart(dst_slide, shape, mapping, src_slide)
    else:
        _migrate_generic_shape(dst_slide, shape, mapping, src_slide)


def _migrate_generic_shape(dst_slide, shape, mapping, src_slide):
    new_el = deepcopy(shape._element)
    new_left, new_top, new_width, new_height = map_position(
        shape.left, shape.top, shape.width, shape.height, mapping
    )
    xfrm = new_el.find('.//a:xfrm', NSMAP)
    if xfrm is None:
        xfrm = new_el.find('p:xfrm', NSMAP)
    if xfrm is not None:
        off = xfrm.find('a:off', NSMAP)
        if off is not None:
            off.set('x', str(new_left))
            off.set('y', str(new_top))
        ext = xfrm.find('a:ext', NSMAP)
        if ext is not None:
            ext.set('cx', str(new_width))
            ext.set('cy', str(new_height))

    r_ns = NSMAP['r']
    r_attrs = [f'{{{r_ns}}}{a}' for a in ('embed', 'link', 'id', 'dm', 'lo', 'qs', 'cs')]
    rId_map = {}
    for el in new_el.iter():
        for attr in r_attrs:
            val = el.get(attr)
            if val and val not in rId_map and val in src_slide.part.rels:
                rel = src_slide.part.rels[val]
                if rel.is_external:
                    new_rId = dst_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[val] = new_rId
    if rId_map:
        _remap_rids(new_el, rId_map)

    dst_slide.shapes._spTree.append(new_el)


def _migrate_chart(dst_slide, shape, mapping, src_slide):
    new_el = deepcopy(shape._element)
    new_left, new_top, new_width, new_height = map_position(
        shape.left, shape.top, shape.width, shape.height, mapping
    )
    xfrm = new_el.find('.//a:xfrm', NSMAP)
    if xfrm is not None:
        off = xfrm.find('a:off', NSMAP)
        if off is not None:
            off.set('x', str(new_left))
            off.set('y', str(new_top))
        ext = xfrm.find('a:ext', NSMAP)
        if ext is not None:
            ext.set('cx', str(new_width))
            ext.set('cy', str(new_height))

    chart_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    chart_ref = new_el.find(f'.//{{{chart_ns}}}chart')
    if chart_ref is not None:
        old_rId = chart_ref.get(f'{{{NSMAP["r"]}}}id')
        if old_rId and old_rId in src_slide.part.rels:
            rel = src_slide.part.rels[old_rId]
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            chart_ref.set(f'{{{NSMAP["r"]}}}id', new_rId)

    dst_slide.shapes._spTree.append(new_el)


def _get_all_shape_ids(slide):
    spids = set()
    for shape in slide.shapes:
        spids.add(str(shape.shape_id))
    spTree = slide._element.find('p:cSld/p:spTree', NSMAP)
    if spTree is not None:
        p_ns = NSMAP['p']
        for mc in spTree.findall('mc:AlternateContent', NSMAP):
            choice = mc.find('mc:Choice', NSMAP)
            if choice is None:
                continue
            for tag_local in SHAPE_TAGS:
                el = choice.find(f'{{{p_ns}}}{tag_local}')
                if el is not None:
                    for cNvPr in el.iter(f'{{{p_ns}}}cNvPr'):
                        sid = cNvPr.get('id')
                        if sid:
                            spids.add(sid)
                        break
    return spids


def _get_content_shape_ids(src_slide):
    spids = {str(s.shape_id) for s in get_content_shapes(src_slide)}
    title_shape = _find_title_shape(src_slide)
    header_bottom = (title_shape.top + title_shape.height + Emu(150000)) if title_shape else 0
    p_ns = NSMAP['p']
    for mc_el, _, _, _, _ in _get_mc_elements(src_slide, header_bottom):
        choice = mc_el.find('mc:Choice', NSMAP)
        if choice is None:
            continue
        for tag_local in SHAPE_TAGS:
            el = choice.find(f'{{{p_ns}}}{tag_local}')
            if el is not None:
                for cNvPr in el.iter(f'{{{p_ns}}}cNvPr'):
                    sid = cNvPr.get('id')
                    if sid:
                        spids.add(sid)
                    break
    return spids


def _deduplicate_template_ids(slide, template_elements):
    p_ns = NSMAP['p']
    spTree = slide._element.find('p:cSld/p:spTree', NSMAP)
    if spTree is None:
        return

    migrated_ids = set()
    for child in spTree:
        if child not in template_elements:
            for cNvPr in child.iter(f'{{{p_ns}}}cNvPr'):
                sid = cNvPr.get('id')
                if sid:
                    migrated_ids.add(sid)
                break

    if not migrated_ids:
        return

    all_ids = set()
    for child in spTree:
        for cNvPr in child.iter(f'{{{p_ns}}}cNvPr'):
            sid = cNvPr.get('id')
            if sid:
                all_ids.add(int(sid))
            break
    next_id = max(all_ids) + 1 if all_ids else 1

    for child in spTree:
        if child not in template_elements:
            continue
        for cNvPr in child.iter(f'{{{p_ns}}}cNvPr'):
            sid = cNvPr.get('id')
            if sid and sid in migrated_ids:
                cNvPr.set('id', str(next_id))
                next_id += 1
            break


def _filter_orphaned_animations(timing_el, valid_spids):
    p_ns = NSMAP['p']
    for seq in timing_el.iter(f'{{{p_ns}}}seq'):
        seq_cTn = seq.find(f'{{{p_ns}}}cTn')
        if seq_cTn is None:
            continue
        childTnLst = seq_cTn.find(f'{{{p_ns}}}childTnLst')
        if childTnLst is None:
            continue
        for click_par in list(childTnLst):
            click_cTn = click_par.find(f'{{{p_ns}}}cTn')
            if click_cTn is None:
                continue
            inner_list = click_cTn.find(f'{{{p_ns}}}childTnLst')
            if inner_list is None:
                continue
            for anim_par in list(inner_list):
                targets = anim_par.findall(f'.//{{{p_ns}}}spTgt')
                if targets and any(t.get('spid') not in valid_spids for t in targets):
                    inner_list.remove(anim_par)
            if len(inner_list) == 0:
                childTnLst.remove(click_par)


def _copy_slide_animations(src_slide, dst_slide, content_spids=None):
    src_el = src_slide._element
    dst_el = dst_slide._element

    for tag in ('p:transition', 'p:timing'):
        dst_node = dst_el.find(tag, NSMAP)
        if dst_node is not None:
            dst_el.remove(dst_node)

        src_node = src_el.find(tag, NSMAP)
        if src_node is None:
            continue
        new_node = deepcopy(src_node)

        r_ns = NSMAP['r']
        r_attrs = [f'{{{r_ns}}}{a}' for a in ('embed', 'link', 'id')]
        rId_map = {}
        for el in new_node.iter():
            for attr in r_attrs:
                val = el.get(attr)
                if val and val not in rId_map and val in src_slide.part.rels:
                    rel = src_slide.part.rels[val]
                    if rel.is_external:
                        new_rId = dst_slide.part.relate_to(
                            rel.target_ref, rel.reltype, is_external=True)
                    else:
                        new_rId = dst_slide.part.relate_to(
                            rel.target_part, rel.reltype)
                    rId_map[val] = new_rId
        if rId_map:
            _remap_rids(new_node, rId_map)

        if tag == 'p:timing':
            valid_spids = content_spids if content_spids is not None else _get_all_shape_ids(dst_slide)
            _filter_orphaned_animations(new_node, valid_spids)

        dst_el.append(new_node)


def restyle(source_path, template_path, output_path):
    if not os.path.exists(source_path):
        raise FileNotFoundError(f"源文件不存在: {source_path}")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")

    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

    template_prs = Presentation(template_path)
    validate_template(template_prs)
    content_bounds = get_content_area_bounds(template_prs.slides[1])

    src_prs = Presentation(source_path)
    if len(src_prs.slides) == 0:
        raise ValueError("源 PPT 没有任何幻灯片")

    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    cover_title = get_title_text(src_prs.slides[0])
    set_title_text(prs.slides[0], cover_title)
    _copy_slide_animations(src_prs.slides[0], prs.slides[0], content_spids=set())

    template_slide_index = 1
    has_content_slides = False

    for i in range(1, len(src_prs.slides)):
        src_slide = src_prs.slides[i]
        title = get_title_text(src_slide)
        shapes = get_content_shapes(src_slide)

        if not shapes and not title:
            continue

        has_content_slides = True
        new_slide = duplicate_slide(prs, template_slide_index)
        set_title_text(new_slide, title)

        spTree = new_slide._element.find('p:cSld/p:spTree', NSMAP)
        template_elements = set()
        for child in spTree:
            tag = _tag_local(child)
            if tag in SHAPE_TAGS or tag == 'AlternateContent':
                template_elements.add(child)

        migrate_content(new_slide, src_slide, content_bounds)
        _deduplicate_template_ids(new_slide, template_elements)

        content_spids = _get_content_shape_ids(src_slide)
        _copy_slide_animations(src_slide, new_slide, content_spids=content_spids)

    remove_slide(prs, template_slide_index)

    if not has_content_slides and len(src_prs.slides) == 1:
        pass  # 仅保留封面

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print("用法: python ppt_restyle.py <源PPT路径> [模板路径] [输出路径]")
        sys.exit(1)

    source = sys.argv[1]
    base_dir = Path(__file__).parent.parent

    template = (
        sys.argv[2] if len(sys.argv) > 2
        else str(base_dir / 'templates' / 'template.pptx')
    )

    if len(sys.argv) > 3:
        output = sys.argv[3]
    else:
        source_name = Path(source).stem
        output = str(base_dir / 'output' / f'{source_name}_styled.pptx')

    try:
        result = restyle(source, template, output)
        print(f"转换完成: {result}")
    except Exception as e:
        print(f"转换失败: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
