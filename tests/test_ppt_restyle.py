import pytest
from pptx import Presentation
from pptx.util import Inches

from scripts.ppt_restyle import validate_template


class TestValidateTemplate:
    def test_valid_template(self, template_pptx):
        prs = Presentation(template_pptx)
        validate_template(prs)

    def test_rejects_single_slide(self, tmp_path):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / "bad.pptx"
        prs.save(str(path))
        with pytest.raises(ValueError, match="2 页"):
            validate_template(Presentation(str(path)))

    def test_rejects_three_slides(self, tmp_path):
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / "bad.pptx"
        prs.save(str(path))
        with pytest.raises(ValueError, match="2 页"):
            validate_template(Presentation(str(path)))


from scripts.ppt_restyle import duplicate_slide, remove_slide


class TestSlideOperations:
    def test_duplicate_slide_adds_one(self, template_pptx):
        prs = Presentation(template_pptx)
        assert len(prs.slides) == 2
        duplicate_slide(prs, 1)
        assert len(prs.slides) == 3

    def test_duplicate_preserves_shapes(self, template_pptx):
        prs = Presentation(template_pptx)
        original = prs.slides[1]
        original_count = len(original.shapes)
        new_slide = duplicate_slide(prs, 1)
        assert len(new_slide.shapes) == original_count

    def test_remove_slide(self, template_pptx):
        prs = Presentation(template_pptx)
        duplicate_slide(prs, 1)
        assert len(prs.slides) == 3
        remove_slide(prs, 2)
        assert len(prs.slides) == 2


from scripts.ppt_restyle import get_title_text, set_title_text


class TestTitleOperations:
    def test_get_title_text(self, template_pptx):
        prs = Presentation(template_pptx)
        assert get_title_text(prs.slides[0]) == "模板标题"

    def test_get_title_text_empty(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        assert get_title_text(slide) == ""

    def test_set_title_text(self, template_pptx):
        prs = Presentation(template_pptx)
        set_title_text(prs.slides[0], "新标题")
        assert get_title_text(prs.slides[0]) == "新标题"


from scripts.ppt_restyle import (
    get_content_shapes, compute_bounding_box,
    compute_mapping, map_position
)


class TestContentAnalysis:
    def test_get_content_shapes_excludes_title(self, source_pptx):
        prs = Presentation(source_pptx)
        slide = prs.slides[1]  # 文本页
        shapes = get_content_shapes(slide)

        def _ph_idx(s):
            try:
                ph = s.placeholder_format
                return ph.idx if ph is not None else None
            except ValueError:
                return None

        titles = [s for s in shapes if get_title_text(slide) and _ph_idx(s) == 0]
        assert len(titles) == 0

    def test_get_content_shapes_includes_textbox(self, source_pptx):
        prs = Presentation(source_pptx)
        slide = prs.slides[1]
        shapes = get_content_shapes(slide)
        texts = [s.text_frame.text for s in shapes if s.has_text_frame]
        assert "额外文本框内容" in texts

    def test_compute_bounding_box(self, source_pptx):
        prs = Presentation(source_pptx)
        shapes = get_content_shapes(prs.slides[1])
        bbox = compute_bounding_box(shapes)
        assert bbox is not None
        left, top, right, bottom = bbox
        assert right > left
        assert bottom > top

    def test_compute_bounding_box_empty(self):
        assert compute_bounding_box([]) is None


class TestCoordinateMapping:
    def test_same_size_scale_one(self):
        src = (0, 0, 1000, 1000)
        dst = (0, 0, 1000, 1000)
        mapping = compute_mapping(src, dst)
        scale = mapping[0]
        assert scale == pytest.approx(1.0)

    def test_half_size_scale(self):
        src = (0, 0, 2000, 2000)
        dst = (0, 0, 1000, 1000)
        mapping = compute_mapping(src, dst)
        scale = mapping[0]
        assert scale == pytest.approx(0.5)

    def test_map_position_identity(self):
        mapping = (1.0, 100, 100, 100, 100)
        result = map_position(200, 200, 50, 50, mapping)
        assert result == (200, 200, 50, 50)

    def test_map_position_with_scale(self):
        mapping = (0.5, 0, 0, 0, 0)
        result = map_position(200, 400, 100, 100, mapping)
        assert result == (100, 200, 50, 50)


from scripts.ppt_restyle import migrate_content, get_content_area_bounds


class TestShapeMigration:
    def test_migrate_textbox(self, template_pptx, source_pptx):
        tpl = Presentation(template_pptx)
        src = Presentation(source_pptx)
        dst_bounds = get_content_area_bounds(tpl.slides[1])

        dst_slide = duplicate_slide(tpl, 1)
        migrate_content(dst_slide, src.slides[1], dst_bounds)

        texts = [s.text_frame.text for s in dst_slide.shapes if s.has_text_frame]
        assert any("额外文本框内容" in t for t in texts)

    def test_migrate_table(self, template_pptx, source_pptx):
        tpl = Presentation(template_pptx)
        src = Presentation(source_pptx)
        dst_bounds = get_content_area_bounds(tpl.slides[1])

        dst_slide = duplicate_slide(tpl, 1)
        migrate_content(dst_slide, src.slides[2], dst_bounds)

        tables = [s for s in dst_slide.shapes if s.has_table]
        assert len(tables) >= 1
        assert tables[0].table.cell(0, 0).text == "A"

    def test_migrate_empty_slide_no_crash(self, template_pptx, tmp_path):
        tpl = Presentation(template_pptx)
        empty_prs = Presentation()
        empty_slide = empty_prs.slides.add_slide(empty_prs.slide_layouts[6])
        dst_bounds = get_content_area_bounds(tpl.slides[1])

        dst_slide = duplicate_slide(tpl, 1)
        migrate_content(dst_slide, empty_slide, dst_bounds)


from scripts.ppt_restyle import restyle


class TestRestyle:
    def test_full_restyle(self, template_pptx, source_pptx, tmp_path):
        output_path = str(tmp_path / "output.pptx")
        restyle(source_pptx, template_pptx, output_path)

        result = Presentation(output_path)
        # 封面 + 2 内容页 = 3 页
        assert len(result.slides) == 3

    def test_cover_title_replaced(self, template_pptx, source_pptx, tmp_path):
        output_path = str(tmp_path / "output.pptx")
        restyle(source_pptx, template_pptx, output_path)

        result = Presentation(output_path)
        assert get_title_text(result.slides[0]) == "我的演示文稿"

    def test_content_title_replaced(self, template_pptx, source_pptx, tmp_path):
        output_path = str(tmp_path / "output.pptx")
        restyle(source_pptx, template_pptx, output_path)

        result = Presentation(output_path)
        assert get_title_text(result.slides[1]) == "第一节标题"
        assert get_title_text(result.slides[2]) == "数据表格"

    def test_single_slide_source(self, template_pptx, source_single_slide_pptx, tmp_path):
        output_path = str(tmp_path / "output.pptx")
        restyle(source_single_slide_pptx, template_pptx, output_path)

        result = Presentation(output_path)
        assert len(result.slides) == 1
        assert get_title_text(result.slides[0]) == "仅封面"

    def test_missing_source_raises(self, template_pptx, tmp_path):
        with pytest.raises(FileNotFoundError):
            restyle("/nonexistent.pptx", template_pptx, str(tmp_path / "out.pptx"))

    def test_missing_template_raises(self, source_pptx, tmp_path):
        with pytest.raises(FileNotFoundError):
            restyle(source_pptx, "/nonexistent.pptx", str(tmp_path / "out.pptx"))


from scripts.ppt_restyle import (
    _copy_slide_animations, _get_content_shape_ids, _filter_orphaned_animations
)
from lxml import etree
from copy import deepcopy


NSMAP_TEST = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def _add_animation_to_slide(slide, shape_ids):
    """Inject a p:timing element targeting the given shape IDs."""
    p_ns = NSMAP_TEST['p']
    timing = etree.SubElement(slide._element, f'{{{p_ns}}}timing')
    tnLst = etree.SubElement(timing, f'{{{p_ns}}}tnLst')
    par = etree.SubElement(tnLst, f'{{{p_ns}}}par')
    cTn = etree.SubElement(par, f'{{{p_ns}}}cTn', attrib={'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'})
    childTnLst = etree.SubElement(cTn, f'{{{p_ns}}}childTnLst')
    seq = etree.SubElement(childTnLst, f'{{{p_ns}}}seq', attrib={'concurrent': '1', 'nextAc': 'seek'})
    seq_cTn = etree.SubElement(seq, f'{{{p_ns}}}cTn', attrib={'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'})
    seq_childTnLst = etree.SubElement(seq_cTn, f'{{{p_ns}}}childTnLst')
    for i, spid in enumerate(shape_ids):
        click_par = etree.SubElement(seq_childTnLst, f'{{{p_ns}}}par')
        click_cTn = etree.SubElement(click_par, f'{{{p_ns}}}cTn', attrib={'id': str(10 + i), 'fill': 'hold'})
        inner_list = etree.SubElement(click_cTn, f'{{{p_ns}}}childTnLst')
        anim_par = etree.SubElement(inner_list, f'{{{p_ns}}}par')
        anim_cTn = etree.SubElement(anim_par, f'{{{p_ns}}}cTn', attrib={'id': str(20 + i), 'presetID': '1', 'presetClass': 'entr', 'fill': 'hold'})
        stCondLst = etree.SubElement(anim_cTn, f'{{{p_ns}}}stCondLst')
        cond = etree.SubElement(stCondLst, f'{{{p_ns}}}cond', attrib={'delay': '0'})
        anim_cTn_child = etree.SubElement(anim_cTn, f'{{{p_ns}}}childTnLst')
        set_el = etree.SubElement(anim_cTn_child, f'{{{p_ns}}}set')
        set_cBhvr = etree.SubElement(set_el, f'{{{p_ns}}}cBhvr')
        set_cTn2 = etree.SubElement(set_cBhvr, f'{{{p_ns}}}cTn', attrib={'id': str(30 + i), 'dur': '1', 'fill': 'hold'})
        tgtEl = etree.SubElement(set_cBhvr, f'{{{p_ns}}}tgtEl')
        spTgt = etree.SubElement(tgtEl, f'{{{p_ns}}}spTgt', attrib={'spid': str(spid)})


class TestAnimationFiltering:
    def test_template_animations_stripped(self, template_pptx, tmp_path):
        """模板自带的动画应在输出中被清除。"""
        tpl_prs = Presentation(template_pptx)
        title_spid = tpl_prs.slides[0].shapes[0].shape_id
        _add_animation_to_slide(tpl_prs.slides[0], [title_spid])
        tpl_path = str(tmp_path / "tpl_with_anim.pptx")
        tpl_prs.save(tpl_path)

        src_prs = Presentation()
        slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])
        slide.placeholders[0].text = "Cover"
        src_path = str(tmp_path / "src_no_anim.pptx")
        src_prs.save(src_path)

        output_path = str(tmp_path / "out.pptx")
        restyle(src_path, tpl_path, output_path)

        result = Presentation(output_path)
        timing = result.slides[0]._element.find(
            'p:timing', NSMAP_TEST
        )
        assert timing is None

    def test_content_animations_kept(self, template_pptx, tmp_path):
        """内容区域形状的动画应保留。"""
        src_prs = Presentation()
        src_prs.slides.add_slide(src_prs.slide_layouts[0])
        src_prs.slides[0].placeholders[0].text = "Cover"
        content_slide = src_prs.slides.add_slide(src_prs.slide_layouts[1])
        content_slide.placeholders[0].text = "Page"
        txBox = content_slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(5), Inches(1)
        )
        txBox.text_frame.text = "animated content"
        content_spid = txBox.shape_id
        _add_animation_to_slide(content_slide, [content_spid])
        src_path = str(tmp_path / "src.pptx")
        src_prs.save(src_path)

        output_path = str(tmp_path / "out.pptx")
        restyle(src_path, template_pptx, output_path)

        result = Presentation(output_path)
        timing = result.slides[1]._element.find('p:timing', NSMAP_TEST)
        assert timing is not None
        targets = timing.findall(
            f'.//{{{NSMAP_TEST["p"]}}}spTgt'
        )
        target_ids = {t.get('spid') for t in targets}
        assert str(content_spid) in target_ids

    def test_title_animations_removed(self, template_pptx, tmp_path):
        """源 PPT 标题形状的动画不应出现在输出中。"""
        src_prs = Presentation()
        src_prs.slides.add_slide(src_prs.slide_layouts[0])
        src_prs.slides[0].placeholders[0].text = "Cover"
        content_slide = src_prs.slides.add_slide(src_prs.slide_layouts[1])
        content_slide.placeholders[0].text = "Titled"
        title_spid = content_slide.placeholders[0].shape_id
        _add_animation_to_slide(content_slide, [title_spid])
        src_path = str(tmp_path / "src.pptx")
        src_prs.save(src_path)

        output_path = str(tmp_path / "out.pptx")
        restyle(src_path, template_pptx, output_path)

        result = Presentation(output_path)
        timing = result.slides[1]._element.find('p:timing', NSMAP_TEST)
        if timing is not None:
            targets = timing.findall(
                f'.//{{{NSMAP_TEST["p"]}}}spTgt'
            )
            target_ids = {t.get('spid') for t in targets}
            assert str(title_spid) not in target_ids

    def test_mixed_animations_filters_correctly(self, template_pptx, tmp_path):
        """同时有标题和内容动画时，只保留内容动画。"""
        src_prs = Presentation()
        src_prs.slides.add_slide(src_prs.slide_layouts[0])
        src_prs.slides[0].placeholders[0].text = "Cover"
        content_slide = src_prs.slides.add_slide(src_prs.slide_layouts[1])
        content_slide.placeholders[0].text = "Mixed"
        txBox = content_slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(5), Inches(1)
        )
        txBox.text_frame.text = "keep this animation"
        title_spid = content_slide.placeholders[0].shape_id
        content_spid = txBox.shape_id
        _add_animation_to_slide(content_slide, [title_spid, content_spid])
        src_path = str(tmp_path / "src.pptx")
        src_prs.save(src_path)

        output_path = str(tmp_path / "out.pptx")
        restyle(src_path, template_pptx, output_path)

        result = Presentation(output_path)
        timing = result.slides[1]._element.find('p:timing', NSMAP_TEST)
        assert timing is not None
        targets = timing.findall(
            f'.//{{{NSMAP_TEST["p"]}}}spTgt'
        )
        target_ids = {t.get('spid') for t in targets}
        assert str(content_spid) in target_ids
        assert str(title_spid) not in target_ids
