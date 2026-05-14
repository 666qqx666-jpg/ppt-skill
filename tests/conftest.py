import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


@pytest.fixture
def template_pptx(tmp_path):
    """生成一个标准模板 PPT：封面页 + 内容页。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 第1页：封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "模板标题"
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

    # 第2页：内容页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = "页面标题"
    slide.placeholders[1].text = "页面内容"
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

    path = tmp_path / "template.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def source_pptx(tmp_path):
    """生成一个包含混合内容的源 PPT：封面 + 文本页 + 表格页。"""
    prs = Presentation()

    # 第1页：封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "我的演示文稿"

    # 第2页：文本内容
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = "第一节标题"
    slide.placeholders[1].text = "这是第一节的内容"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
    txBox.text_frame.text = "额外文本框内容"

    # 第3页：表格
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = "数据表格"
    table_shape = slide.shapes.add_table(
        2, 3, Inches(1), Inches(2), Inches(8), Inches(2)
    )
    tbl = table_shape.table
    for row_idx, row_data in enumerate([["A", "B", "C"], ["1", "2", "3"]]):
        for col_idx, val in enumerate(row_data):
            tbl.cell(row_idx, col_idx).text = val

    path = tmp_path / "source.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def source_single_slide_pptx(tmp_path):
    """仅有封面的源 PPT。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "仅封面"
    path = tmp_path / "single.pptx"
    prs.save(str(path))
    return str(path)
